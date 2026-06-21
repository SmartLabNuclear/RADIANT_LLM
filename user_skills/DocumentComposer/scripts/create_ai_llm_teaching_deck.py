"""
Create a PowerPoint teaching deck: AI, Generative AI, LLMs, RAG, CAG, MCP, and Agent Skills.

Output:
    AI_LLMs_Generative_AI_High_School_Teaching_Deck_recreated.pptx

Dependencies:
    pip install python-pptx pillow

Usage:
    python create_ai_llm_teaching_deck.py

Notes:
    - The script uses local illustration files if they are found next to this script or in /mnt/data.
    - If images are not available, it still generates the full deck using shapes and diagrams.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

try:
    from PIL import Image
except Exception:  # pragma: no cover
    Image = None


# -----------------------------
# Theme constants
# -----------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(18, 32, 58)
BLUE = RGBColor(43, 111, 246)
CYAN = RGBColor(0, 184, 217)
PURPLE = RGBColor(124, 77, 255)
GREEN = RGBColor(25, 170, 112)
ORANGE = RGBColor(245, 153, 61)
RED = RGBColor(226, 80, 86)
YELLOW = RGBColor(250, 201, 80)
WHITE = RGBColor(255, 255, 255)
OFFWHITE = RGBColor(247, 250, 255)
LIGHT = RGBColor(226, 234, 247)
MID = RGBColor(91, 110, 140)
DARK = RGBColor(32, 45, 71)
GRAY = RGBColor(128, 143, 166)
PALE_BLUE = RGBColor(233, 240, 255)
PALE_PURPLE = RGBColor(242, 236, 255)
PALE_GREEN = RGBColor(232, 248, 241)
PALE_ORANGE = RGBColor(255, 244, 230)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"

SCRIPT_DIR = Path(__file__).resolve().parent

IMAGE_CANDIDATES = {
    "classroom": [
        "a_bright_modern_classroom_learning_space_scene_in.png",
        "a_bright_modern_classroom_learning_space_scene_in_batch_1.png",
    ],
    "generative": [
        "a_wide_clean_colorful_infographic_concept_illu.png",
        "a_wide_clean_colorful_infographic_concept_illu_batch_2.png",
    ],
    "rag": [
        "a_clean_high_quality_3d_cgi_illustration_with_a_s.png",
        "a_clean_high_quality_3d_cgi_illustration_with_a_s_batch_3.png",
    ],
    "agents": [
        "wide_clean_3d_illustration_style_scene_in_a_softl.png",
        "wide_clean_3d_illustration_style_scene_in_a_softl_batch_4.png",
    ],
}


def find_asset(kind: str) -> Path | None:
    """Find an optional image asset by name."""
    for name in IMAGE_CANDIDATES.get(kind, []):
        for folder in (SCRIPT_DIR, Path("/mnt/data"), Path.cwd()):
            candidate = folder / name
            if candidate.exists():
                return candidate
    return None


# -----------------------------
# Basic drawing helpers
# -----------------------------
def rgb(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def set_fill(shape, color: RGBColor, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor = LIGHT, width: float = 1.0, transparency: float = 0.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 24,
    color: RGBColor = DARK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = FONT_BODY,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    items: Sequence[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 21,
    color: RGBColor = DARK,
    bullet_color: RGBColor | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = FONT_BODY
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.line_spacing = 1.05
        # python-pptx has limited public bullet formatting. Use a bullet character for reliable rendering.
        p.text = f"• {item}"
        if bullet_color:
            p.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    color = WHITE if dark else NAVY
    add_text(slide, title, 0.65, 0.38, 12.0, 0.55, size=30, color=color, bold=True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, 0.67, 0.98, 11.7, 0.38, size=15, color=LIGHT if dark else MID)


def add_footer(slide, n: int) -> None:
    add_text(slide, f"AI & LLMs for High School • {n}", 10.15, 7.05, 2.75, 0.22, size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def add_background(slide, color: RGBColor = OFFWHITE) -> None:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(rect, color)
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_top_accent(slide, color1: RGBColor = BLUE, color2: RGBColor = CYAN) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    set_fill(bar, color1)
    bar.line.fill.background()
    small = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.2), 0, Inches(4.2), Inches(0.12))
    set_fill(small, color2)
    small.line.fill.background()


def add_round_rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor = LIGHT, radius=True, transparency=0.0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, fill, transparency)
    set_line(shp, line, 1)
    return shp


def add_pill(slide, text: str, x, y, w, h, fill: RGBColor, text_color: RGBColor = WHITE, size: int = 14):
    shp = add_round_rect(slide, x, y, w, h, fill, fill)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp


def add_card(slide, title: str, body: str, x, y, w, h, fill: RGBColor = WHITE, accent: RGBColor = BLUE, icon: str | None = None):
    card = add_round_rect(slide, x, y, w, h, fill, LIGHT)
    add_round_rect(slide, x, y, 0.11, h, accent, accent, radius=False)
    if icon:
        add_text(slide, icon, x + 0.23, y + 0.18, 0.42, 0.3, size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
        tx = x + 0.72
        tw = w - 0.9
    else:
        tx = x + 0.25
        tw = w - 0.45
    add_text(slide, title, tx, y + 0.18, tw, 0.32, size=16, color=NAVY, bold=True)
    add_text(slide, body, tx, y + 0.58, tw, h - 0.72, size=11, color=MID)
    return card


def add_arrow(slide, x1, y1, x2, y2, color: RGBColor = BLUE, width: float = 2.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_image_cover(slide, image_path: Path, x, y, w, h) -> bool:
    """Add image cropped to cover a rectangle. Returns False if image cannot be loaded."""
    if Image is None or image_path is None or not image_path.exists():
        return False
    try:
        with Image.open(image_path) as im:
            img_w, img_h = im.size
        box_ratio = w / h
        img_ratio = img_w / img_h
        pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        if img_ratio > box_ratio:
            crop = (img_ratio - box_ratio) / (2 * img_ratio)
            pic.crop_left = crop
            pic.crop_right = crop
        else:
            crop = (box_ratio - img_ratio) / (2 * box_ratio)
            pic.crop_top = crop
            pic.crop_bottom = crop
        return True
    except Exception:
        return False


def add_section_label(slide, text: str, color: RGBColor = BLUE):
    add_pill(slide, text.upper(), 0.65, 6.75, 2.25, 0.33, color, WHITE, size=10)


# -----------------------------
# Diagram helpers
# -----------------------------
def add_vocabulary_stack(slide, x, y):
    levels = [
        ("AI", "Machines doing intelligent tasks", BLUE, 5.3),
        ("Machine Learning", "Learns patterns from data", CYAN, 4.45),
        ("Deep Learning", "Uses neural networks", PURPLE, 3.6),
        ("Generative AI", "Creates new content", GREEN, 2.75),
        ("LLMs", "Language-focused generative models", ORANGE, 1.9),
    ]
    for i, (name, desc, col, width) in enumerate(levels):
        y_i = y + i * 0.78
        add_round_rect(slide, x + (5.3 - width) / 2, y_i, width, 0.58, WHITE, col)
        add_text(slide, name, x + (5.3 - width) / 2 + 0.12, y_i + 0.08, width - 0.24, 0.2, size=14, color=col, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + (5.3 - width) / 2 + 0.12, y_i + 0.32, width - 0.24, 0.18, size=8, color=MID, align=PP_ALIGN.CENTER)


def add_token_predictor(slide):
    sentence = ["The", "quick", "brown", "fox", "jumps", "over", "the", "___"]
    x = 1.0
    for i, word in enumerate(sentence):
        fill = PALE_BLUE if word != "___" else PALE_ORANGE
        accent = BLUE if word != "___" else ORANGE
        add_round_rect(slide, x + i * 1.22, 2.55, 1.0, 0.55, fill, accent)
        add_text(slide, word, x + i * 1.22 + 0.05, 2.72, 0.9, 0.16, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    options = [("fence", 0.46), ("moon", 0.18), ("algorithm", 0.02)]
    for i, (tok, prob) in enumerate(options):
        y = 4.0 + i * 0.55
        add_text(slide, tok, 2.4, y, 1.3, 0.25, size=15, color=NAVY, bold=True)
        add_round_rect(slide, 4.0, y + 0.03, 5.8 * prob / 0.46, 0.18, BLUE if i == 0 else CYAN if i == 1 else LIGHT, BLUE if i == 0 else CYAN if i == 1 else LIGHT)
        add_text(slide, f"{int(prob*100)}%", 10.0, y, 0.6, 0.22, size=11, color=MID)


def add_process_row(slide, labels: Sequence[tuple[str, str, RGBColor]], x, y, card_w=2.15, card_h=1.05, gap=0.36):
    for i, (title, body, color) in enumerate(labels):
        cx = x + i * (card_w + gap)
        add_round_rect(slide, cx, y, card_w, card_h, WHITE, color)
        add_text(slide, title, cx + 0.12, y + 0.15, card_w - 0.24, 0.25, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, cx + 0.15, y + 0.48, card_w - 0.3, 0.38, size=9, color=MID, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_arrow(slide, cx + card_w + 0.04, y + card_h / 2, cx + card_w + gap - 0.08, y + card_h / 2, MID, 1.5)


def add_clear_checklist(slide):
    items = [
        ("C", "Context", "Who is the answer for? What background matters?", BLUE),
        ("L", "Limits", "Length, format, grade level, allowed sources.", CYAN),
        ("E", "Examples", "Show the style or type of answer you want.", PURPLE),
        ("A", "Ask", "State the exact task or question clearly.", GREEN),
        ("R", "Review", "Check the answer, then ask for improvements.", ORANGE),
    ]
    for i, (letter, name, desc, color) in enumerate(items):
        y = 1.7 + i * 0.84
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.0), Inches(y), Inches(0.52), Inches(0.52))
        set_fill(circ, color)
        circ.line.fill.background()
        add_text(slide, letter, 1.0, y + 0.09, 0.52, 0.18, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, name, 1.75, y, 1.45, 0.26, size=16, color=NAVY, bold=True)
        add_text(slide, desc, 3.15, y, 8.4, 0.3, size=14, color=MID)


# -----------------------------
# Slide builders
# -----------------------------
def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    def new_slide(bg: RGBColor = OFFWHITE):
        s = prs.slides.add_slide(blank)
        add_background(s, bg)
        add_top_accent(s)
        return s

    # 1 Title
    slide = prs.slides.add_slide(blank)
    add_background(slide, NAVY)
    img = find_asset("classroom")
    if img:
        add_image_cover(slide, img, 6.55, 0, 6.78, 7.5)
        overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.55), 0, Inches(6.78), SLIDE_H)
        set_fill(overlay, NAVY, 55)
        overlay.line.fill.background()
    add_text(slide, "AI & Large Language Models", 0.75, 1.05, 5.8, 1.2, size=39, color=WHITE, bold=True, font=FONT_HEAD)
    add_text(slide, "A visual teaching deck for high school students", 0.8, 2.35, 5.2, 0.4, size=18, color=LIGHT)
    add_pill(slide, "Generative AI", 0.8, 3.25, 1.65, 0.38, BLUE)
    add_pill(slide, "LLMs", 2.65, 3.25, 1.0, 0.38, PURPLE)
    add_pill(slide, "RAG / CAG", 3.85, 3.25, 1.35, 0.38, GREEN)
    add_pill(slide, "Agents", 5.4, 3.25, 1.0, 0.38, ORANGE)
    add_text(slide, "Teaching focus: what these systems are, how they work, how to use them, and how to stay safe.", 0.8, 5.7, 5.6, 0.65, size=16, color=LIGHT)

    # 2 Learning goals
    slide = new_slide()
    add_title(slide, "Learning goals", "By the end, students should be able to explain what AI systems do and use them responsibly.")
    goals = [
        ("1", "Define AI, Generative AI, and LLMs", "Vocabulary for the rest of the unit", BLUE),
        ("2", "Describe how LLMs generate text", "Tokens, training, embeddings, and attention", CYAN),
        ("3", "Use prompt engineering", "Give clear instructions, context, examples, and constraints", PURPLE),
        ("4", "Compare RAG, CAG, MCP, and agents", "Understand modern AI app patterns", GREEN),
        ("5", "Evaluate risks and responsible use", "Bias, hallucinations, privacy, cheating, and verification", ORANGE),
    ]
    for i, (num, title, body, color) in enumerate(goals):
        x = 0.9 + (i % 2) * 5.8
        y = 1.75 + (i // 2) * 1.55
        add_card(slide, title, body, x, y, 5.25, 1.1, WHITE, color, icon=num)
    add_footer(slide, 2)

    # 3 What is AI?
    slide = new_slide()
    add_title(slide, "What is AI?", "AI is software that performs tasks that usually require human intelligence.")
    add_bullets(slide, [
        "Recognizes patterns in data",
        "Makes predictions or decisions",
        "Learns from examples instead of fixed step-by-step rules",
        "Can be narrow: one task, not a human mind",
    ], 0.9, 1.7, 6.0, 2.4, size=20)
    add_round_rect(slide, 7.2, 1.65, 4.85, 3.0, PALE_BLUE, BLUE)
    add_text(slide, "AI is not magic", 7.55, 2.0, 4.2, 0.35, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "It is math + data + computing power + human design.", 7.65, 2.7, 4.0, 0.75, size=21, color=DARK, align=PP_ALIGN.CENTER)
    add_section_label(slide, "Big idea", BLUE)
    add_footer(slide, 3)

    # 4 Vocabulary stack
    slide = new_slide()
    add_title(slide, "AI, ML, Deep Learning, Generative AI", "These terms are related, but not the same.")
    add_vocabulary_stack(slide, 0.9, 1.55)
    add_card(slide, "Example", "A face-unlock system may use deep learning to classify whether an image matches the phone owner.", 7.0, 1.7, 4.8, 1.15, WHITE, BLUE)
    add_card(slide, "Generative example", "A text-to-image system creates a new picture from a written prompt.", 7.0, 3.25, 4.8, 1.15, WHITE, GREEN)
    add_card(slide, "LLM example", "A chatbot writes, explains, summarizes, translates, or helps code using language patterns.", 7.0, 4.8, 4.8, 1.15, WHITE, ORANGE)
    add_footer(slide, 4)

    # 5 Generative AI
    slide = new_slide()
    add_title(slide, "What makes AI “generative”?", "It produces new content rather than only classifying or predicting labels.")
    img = find_asset("generative")
    if img:
        add_image_cover(slide, img, 0.7, 1.55, 5.7, 4.0)
    else:
        add_round_rect(slide, 0.7, 1.55, 5.7, 4.0, PALE_PURPLE, PURPLE)
        add_text(slide, "Data → AI model → New content", 1.15, 3.35, 4.8, 0.3, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    content_types = [
        ("Text", "essays, summaries, tutoring", BLUE),
        ("Images", "art, diagrams, storyboards", PURPLE),
        ("Code", "programs, debugging help", CYAN),
        ("Audio / Video", "voices, music, clips", ORANGE),
    ]
    for i, (title, body, color) in enumerate(content_types):
        add_card(slide, title, body, 6.9 + (i % 2) * 2.75, 1.7 + (i // 2) * 1.55, 2.35, 1.1, WHITE, color)
    add_text(slide, "Generative AI learns patterns from examples, then creates new combinations that fit those patterns.", 7.0, 5.45, 5.25, 0.55, size=16, color=DARK, bold=True)
    add_footer(slide, 5)

    # 6 LLM core
    slide = new_slide()
    add_title(slide, "Large Language Models: the core idea", "An LLM is trained to predict likely next tokens from previous tokens.")
    add_token_predictor(slide)
    add_card(slide, "Token", "A piece of text such as a word, part of a word, punctuation mark, or symbol.", 1.0, 5.55, 3.55, 0.95, WHITE, BLUE)
    add_card(slide, "Generation", "The model chooses one token, adds it to the context, then repeats.", 4.9, 5.55, 3.55, 0.95, WHITE, GREEN)
    add_card(slide, "Important", "It does not copy a perfect answer from a database; it predicts likely text from patterns.", 8.8, 5.55, 3.2, 0.95, WHITE, ORANGE)
    add_footer(slide, 6)

    # 7 Inside LLM
    slide = new_slide()
    add_title(slide, "Inside an LLM: tokens, embeddings, attention", "LLMs turn words into numbers, compare context, then generate the next token.")
    steps = [
        ("Tokenize", "Split text into chunks", BLUE),
        ("Embed", "Convert chunks into vectors", CYAN),
        ("Attend", "Focus on relevant context", PURPLE),
        ("Predict", "Choose likely next token", GREEN),
        ("Repeat", "Continue until complete", ORANGE),
    ]
    for i, (title, body, color) in enumerate(steps):
        add_card(slide, title, body, 0.8 + i * 2.45, 2.25, 2.0, 1.25, WHITE, color, icon=str(i + 1))
        if i < 4:
            add_arrow(slide, 2.75 + i * 2.45, 2.86, 3.14 + i * 2.45, 2.86, MID, 1.4)
    add_round_rect(slide, 1.1, 4.75, 10.9, 1.15, PALE_BLUE, BLUE)
    add_text(slide, "Attention helps the model connect words that matter to each other, even when they are far apart in the sentence.", 1.45, 5.1, 10.2, 0.32, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8 Training
    slide = new_slide()
    add_title(slide, "How models are trained", "Training teaches the model patterns; alignment teaches it to be helpful and safer.")
    labels = [
        ("Data", "Books, websites, code, images, examples", BLUE),
        ("Pretraining", "Predict missing/next tokens at huge scale", CYAN),
        ("Fine-tuning", "Improve behavior on task examples", PURPLE),
        ("Feedback", "Humans or AI rank better answers", GREEN),
        ("Deployment", "Model is placed inside an app", ORANGE),
    ]
    add_process_row(slide, labels, 0.55, 2.0, card_w=2.2, card_h=1.25, gap=0.28)
    add_card(slide, "Analogy", "Training is like reading millions of practice problems and learning patterns. It is not the same as understanding the world like a human.", 1.05, 4.45, 5.2, 1.1, WHITE, BLUE)
    add_card(slide, "Compute", "Training large models requires powerful computer chips, energy, engineering, and evaluation.", 7.05, 4.45, 5.2, 1.1, WHITE, ORANGE)
    add_footer(slide, 8)

    # 9 Capabilities
    slide = new_slide()
    add_title(slide, "What can LLMs do well?", "LLMs are especially useful when language, ideas, and structure matter.")
    cards = [
        ("Tutor", "Explain steps, ask questions, give hints", BLUE),
        ("Writer", "Draft, revise, summarize, translate", PURPLE),
        ("Coder", "Explain code, find bugs, build examples", CYAN),
        ("Research helper", "Organize notes and compare sources", GREEN),
        ("Creative partner", "Brainstorm stories, images, games, plans", ORANGE),
        ("Accessibility aid", "Simplify text and support communication", RED),
    ]
    for i, (title, body, color) in enumerate(cards):
        add_card(slide, title, body, 0.8 + (i % 3) * 4.1, 1.75 + (i // 3) * 1.75, 3.55, 1.2, WHITE, color)
    add_text(slide, "Strong use case: helping students think through a task. Weak use case: blindly copying an answer.", 1.0, 5.75, 11.0, 0.38, size=19, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10 Limits
    slide = new_slide()
    add_title(slide, "Limits: why LLMs can be wrong", "They predict plausible text, so they can sound confident even when facts are false.")
    add_bullets(slide, [
        "Hallucination: invented facts, sources, or details",
        "Bias: patterns from data can reflect society’s unfairness",
        "Stale knowledge: training data can be out of date",
        "No true memory unless a system gives it memory",
        "Privacy risk: never paste secrets or personal data",
    ], 0.9, 1.55, 6.2, 3.6, size=18)
    add_round_rect(slide, 7.4, 1.65, 4.55, 3.75, PALE_ORANGE, ORANGE)
    add_text(slide, "Healthy habit", 7.85, 2.05, 3.65, 0.4, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Trust, but verify.", 7.95, 2.9, 3.45, 0.5, size=30, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Use reliable sources, compare answers, and explain the reasoning in your own words.", 7.85, 4.05, 3.65, 0.7, size=15, color=DARK, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)

    # 11 Prompt Engineering
    slide = new_slide()
    add_title(slide, "Prompt Engineering: how to ask better questions", "A prompt is the instruction you give the AI. Better prompts make better outputs.")
    items = [
        ("Role", "“Act as a biology tutor…”", BLUE),
        ("Task", "“Explain photosynthesis…”", CYAN),
        ("Context", "“For a 9th grade class…”", PURPLE),
        ("Format", "“Use a table and 3 examples…”", GREEN),
        ("Constraints", "“Keep it under 150 words…”", ORANGE),
        ("Check", "“List assumptions and uncertainties…”", RED),
    ]
    for i, (title, body, color) in enumerate(items):
        add_card(slide, title, body, 0.8 + (i % 3) * 4.1, 1.75 + (i // 3) * 1.55, 3.55, 1.08, WHITE, color)
    add_round_rect(slide, 1.2, 5.45, 10.8, 0.75, PALE_BLUE, BLUE)
    add_text(slide, "Good prompt = clear goal + useful context + output format + request to check uncertainty", 1.45, 5.68, 10.3, 0.28, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12 CLEAR pattern
    slide = new_slide()
    add_title(slide, "Prompt pattern: CLEAR", "A simple checklist students can remember.")
    add_clear_checklist(slide)
    add_round_rect(slide, 8.3, 1.85, 3.65, 3.35, PALE_GREEN, GREEN)
    add_text(slide, "Example prompt", 8.65, 2.15, 2.95, 0.3, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Act as a patient algebra tutor. Explain how to solve 2x + 5 = 17 for a 9th grader. Show each step, then give one practice problem. Keep it friendly.", 8.65, 2.75, 2.95, 1.75, size=15, color=DARK, align=PP_ALIGN.CENTER)
    add_footer(slide, 12)

    # 13 RAG
    slide = new_slide()
    add_title(slide, "RAG: Retrieval-Augmented Generation", "RAG lets an AI look up relevant information before answering.")
    img = find_asset("rag")
    if img:
        add_image_cover(slide, img, 0.75, 1.6, 5.3, 3.2)
    else:
        add_round_rect(slide, 0.75, 1.6, 5.3, 3.2, PALE_BLUE, BLUE)
        add_text(slide, "Question → Search → Evidence → Answer", 1.05, 3.05, 4.7, 0.3, size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_process_row(slide, [
        ("Question", "Student asks", BLUE),
        ("Search", "Find relevant chunks", CYAN),
        ("Prompt", "Add evidence", PURPLE),
        ("Answer", "Respond with sources", GREEN),
    ], 6.45, 2.0, card_w=1.45, card_h=1.1, gap=0.25)
    add_text(slide, "Question → Search documents → Add evidence to prompt → Generate answer with sources", 6.7, 4.55, 5.1, 0.48, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, "Why it matters", "RAG helps an AI answer from class notes, current documents, or a trusted knowledge base instead of relying only on training data.", 1.15, 5.35, 10.8, 0.9, WHITE, BLUE)
    add_footer(slide, 13)

    # 14 RAG pipeline
    slide = new_slide()
    add_title(slide, "RAG pipeline: one practical view", "A school chatbot could answer from class notes, a textbook, and teacher policies.")
    labels = [
        ("Upload content", "PDFs, websites, notes", BLUE),
        ("Chunk & embed", "Make searchable vectors", CYAN),
        ("Store", "Vector database", PURPLE),
        ("Retrieve", "Find top matches", GREEN),
        ("Generate", "Answer with evidence", ORANGE),
    ]
    add_process_row(slide, labels, 0.55, 1.95, card_w=2.15, card_h=1.2, gap=0.3)
    add_card(slide, "Vector database", "A special index that stores numeric representations of text so similar meanings can be found quickly.", 0.9, 4.25, 3.7, 1.25, WHITE, PURPLE)
    add_card(slide, "Citation habit", "A strong RAG system shows where information came from so users can verify it.", 4.85, 4.25, 3.7, 1.25, WHITE, GREEN)
    add_card(slide, "Failure mode", "Bad retrieval can produce bad answers, even if the model sounds polished.", 8.8, 4.25, 3.7, 1.25, WHITE, RED)
    add_footer(slide, 14)

    # 15 CAG
    slide = new_slide()
    add_title(slide, "CAG: Cache-Augmented Generation", "CAG preloads stable knowledge and reuses cached context for faster answers.")
    add_card(slide, "RAG", "• Searches at question time\n• Good for large or changing collections\n• Can cite retrieved passages\n• More moving parts and possible retrieval mistakes", 0.95, 1.75, 5.35, 3.55, WHITE, BLUE)
    add_card(slide, "CAG", "• Preloads or caches important context\n• Good for stable, repeatedly used knowledge\n• Can be faster and simpler for small/medium knowledge bases\n• Cache must be refreshed when content changes", 7.0, 1.75, 5.35, 3.55, WHITE, GREEN)
    add_text(slide, "Simple analogy: RAG is like searching the library each time. CAG is like keeping the most useful pages already open on your desk.", 1.0, 5.75, 11.2, 0.42, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 15)

    # 16 MCP
    slide = new_slide()
    add_title(slide, "MCP: Model Context Protocol", "MCP is a common way for AI apps to connect with tools and data sources.")
    add_round_rect(slide, 0.95, 1.65, 3.25, 3.4, PALE_BLUE, BLUE)
    add_text(slide, "AI app", 1.55, 2.95, 2.05, 0.32, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.35, 3.35, 5.7, 3.35, BLUE)
    add_round_rect(slide, 5.75, 2.35, 2.0, 1.95, PALE_PURPLE, PURPLE)
    add_text(slide, "MCP\nserver", 6.05, 2.95, 1.4, 0.55, size=22, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 7.85, 3.35, 9.15, 3.35, BLUE)
    tools = [("Files", 9.3, 1.6, BLUE), ("Database", 10.65, 2.55, GREEN), ("Calendar", 9.3, 3.5, ORANGE), ("Web/API", 10.65, 4.45, CYAN)]
    for name, x, y, color in tools:
        add_round_rect(slide, x, y, 1.25, 0.72, WHITE, color)
        add_text(slide, name, x + 0.05, y + 0.21, 1.15, 0.18, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "Think of MCP like a universal adapter",
        "It lets models access approved tools through standardized servers",
        "Permissions and safety rules matter",
    ], 1.0, 5.35, 11.0, 0.95, size=15)
    add_footer(slide, 16)

    # 17 Agents
    slide = new_slide()
    add_title(slide, "Agents and Agent Skills", "Agents are AI systems that can plan, use tools, and complete multi-step tasks.")
    img = find_asset("agents")
    if img:
        add_image_cover(slide, img, 7.15, 1.55, 4.9, 3.3)
    else:
        add_round_rect(slide, 7.15, 1.55, 4.9, 3.3, PALE_BLUE, BLUE)
    steps = [
        ("Goal", "What needs to be done?", BLUE),
        ("Plan", "Break into steps", CYAN),
        ("Tools", "Search, code, calendar, files", PURPLE),
        ("Check", "Evaluate result", GREEN),
        ("Deliver", "Return useful output", ORANGE),
    ]
    for i, (title, body, color) in enumerate(steps):
        add_card(slide, title, body, 0.9, 1.55 + i * 0.85, 5.1, 0.68, WHITE, color)
    add_card(slide, "Agent Skills", "Reusable instructions, code, and workflows that help an agent perform a specialized task consistently.", 1.0, 6.0, 11.0, 0.75, WHITE, PURPLE)
    add_footer(slide, 17)

    # 18 Responsible use
    slide = new_slide()
    add_title(slide, "Responsible use in school", "AI can support learning, but it should not replace thinking.")
    good = [
        ("Use AI as a tutor", "Ask for hints, examples, and feedback", GREEN),
        ("Show your thinking", "Explain how you got your answer", BLUE),
        ("Protect privacy", "Do not paste personal data or secrets", ORANGE),
        ("Check facts", "Compare with trusted sources", PURPLE),
    ]
    bad = [
        ("Do not copy blindly", "That is not learning", RED),
        ("Do not fake sources", "Citations must be real", RED),
        ("Do not bypass rules", "Follow teacher and school policy", RED),
    ]
    for i, (title, body, color) in enumerate(good):
        add_card(slide, title, body, 0.9, 1.65 + i * 1.05, 5.4, 0.78, WHITE, color)
    for i, (title, body, color) in enumerate(bad):
        add_card(slide, title, body, 7.15, 1.9 + i * 1.25, 4.85, 0.85, WHITE, color)
    add_text(slide, "Best rule: AI may help you learn, but your submitted work should reflect your understanding.", 1.05, 6.35, 11.0, 0.3, size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 18)

    # 19 Activity
    slide = new_slide()
    add_title(slide, "Class activity: build a mini RAG answer", "Students practice asking, retrieving, answering, and checking.")
    activity = [
        ("1. Question", "Choose a class topic question", BLUE),
        ("2. Evidence", "Find two reliable passages", CYAN),
        ("3. Prompt", "Ask AI to answer using only the evidence", PURPLE),
        ("4. Verify", "Mark each claim as supported or unsupported", GREEN),
        ("5. Improve", "Revise the answer and add citations", ORANGE),
    ]
    add_process_row(slide, activity, 0.55, 2.0, card_w=2.15, card_h=1.3, gap=0.3)
    add_round_rect(slide, 1.0, 4.35, 11.0, 1.25, PALE_BLUE, BLUE)
    add_text(slide, "Discussion question", 1.35, 4.62, 2.3, 0.25, size=18, color=BLUE, bold=True)
    add_text(slide, "How did the answer change when the AI had evidence compared with when it answered from memory?", 3.55, 4.58, 7.9, 0.34, size=18, color=NAVY, bold=True)
    add_footer(slide, 19)

    # 20 Takeaways
    slide = new_slide()
    add_title(slide, "Key takeaways", "AI is powerful, but students should learn how it works and how to use it responsibly.")
    add_bullets(slide, [
        "LLMs generate by predicting tokens using patterns learned from data",
        "Generative AI can create text, images, code, audio, and more",
        "Prompt engineering improves instructions and outputs",
        "RAG adds external evidence; CAG caches stable knowledge for speed",
        "MCP and agents connect models to tools, but need permissions and guardrails",
        "Best practice: use AI to learn, then verify and explain",
    ], 0.95, 1.55, 8.0, 4.15, size=18)
    add_round_rect(slide, 9.35, 2.05, 2.75, 2.75, PALE_GREEN, GREEN)
    add_text(slide, "Think first.\nAsk well.\nCheck carefully.", 9.65, 2.65, 2.1, 1.3, size=25, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 20)

    # 21 Sources
    slide = new_slide()
    add_title(slide, "Sources and further reading", "Teacher reference links used to prepare this deck.")
    sources = [
        "OpenAI Developers: Agents SDK and tool orchestration examples",
        "Microsoft Learn: Prompt engineering concepts",
        "Model Context Protocol official documentation",
        "Anthropic: Introducing the Model Context Protocol",
        "IBM Developer: Optimizing LLMs with cache augmented generation",
        "Chan et al. (2024): Don't Do RAG: When Cache-Augmented Generation is All You Need",
    ]
    add_bullets(slide, sources, 0.95, 1.55, 11.2, 3.8, size=17)
    add_card(slide, "Suggested teacher note", "Ask students to compare AI output with at least two reliable sources when facts matter.", 1.0, 5.85, 11.0, 0.75, WHITE, BLUE)
    add_footer(slide, 21)

    prs.save(output)


if __name__ == "__main__":
    out = SCRIPT_DIR / "AI_LLMs_Generative_AI_High_School_Teaching_Deck_recreated.pptx"
    build_deck(out)
    print(f"Saved presentation to: {out}")
