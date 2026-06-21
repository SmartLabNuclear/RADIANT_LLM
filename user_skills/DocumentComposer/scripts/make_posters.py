
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyArrowPatch, Circle

wd = Path('/workdir/Articles')
outdir = wd
imgdir = wd / 'poster_assets'
imgdir.mkdir(exist_ok=True)

navy = '#143E6E'; teal = '#1B9AAA'; red = '#C0392B'; gold = '#D4A017'; green = '#2E8B57'; gray = '#5B6573'; light = '#F6F8FB'

def savefig(path):
    plt.tight_layout()
    plt.savefig(path, dpi=220, bbox_inches='tight', facecolor='white')
    plt.close()

def fig_a1_workflow(path):
    fig, ax = plt.subplots(figsize=(7,4.2))
    ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis('off')
    boxes = [(0.3, 3.8, 2.0, 1.2, 'Previously identified\nINH-resistant MTB\nisolates'),(2.8, 3.8, 2.3, 1.2, 'Subset lacking canonical\nkatG/inhA promoter\nmutations (13.7%)'),(5.6, 3.8, 1.7, 1.2, 'MIC testing\n(resazurin assay)'),(7.8, 3.8, 1.8, 1.2, 'WGS\nPacBio Onso')]
    for x,y,w,h,label in boxes:
        ax.add_patch(Rectangle((x,y),w,h,facecolor=light,edgecolor=navy,linewidth=2))
        ax.text(x+w/2,y+h/2,label,ha='center',va='center',fontsize=11,color=navy,weight='bold')
    for i in range(len(boxes)-1):
        x1 = boxes[i][0]+boxes[i][2]; x2 = boxes[i+1][0]
        ax.add_patch(FancyArrowPatch((x1+0.05,4.4),(x2-0.05,4.4),arrowstyle='-|>',mutation_scale=16,color=teal,linewidth=2))
    ax.add_patch(Rectangle((2.1, 1.0), 5.8, 1.5, facecolor='#EAF5FF', edgecolor=teal, linewidth=2))
    ax.text(5.0, 1.75, 'Integrated analysis:\nvariant calling with TB-Profiler v6.6.2 + lineage assignment + MIC comparison', ha='center', va='center', fontsize=11, color=navy)
    ax.set_title('Study workflow', fontsize=15, color=navy, weight='bold')
    savefig(path)

def fig_a1_variant_classes(path):
    fig, ax = plt.subplots(figsize=(7,4.2))
    ax.set_xlim(0, 2); ax.set_ylim(-0.5, 3.8)
    ax.set_xticks([0.5,1.5]); ax.set_xticklabels(['Low-level\nresistance','High-level\nresistance'])
    ax.set_yticks(range(4)); ax.set_yticklabels(['ahpC c.-74G>A','inhA c.-154G>A','katG p.Ser315Asn','katG p.Gly630fs'])
    ax.invert_yaxis(); vals = [0,0,1,1]; colors=[gold,gold,red,red]
    for i,v in enumerate(vals):
        ax.add_patch(Rectangle((v+0.15,i-0.35),0.7,0.7,facecolor=colors[i],edgecolor='white'))
    ax.text(0.5,3.45,'Borderline/low MIC\n(~0.25 µg/mL)',ha='center',va='top',fontsize=10,color=gray)
    ax.text(1.5,3.45,'High MIC\n(≥8 µg/mL)',ha='center',va='top',fontsize=10,color=gray)
    ax.set_title('Detected resistance-associated variants and resistance class', fontsize=14, color=navy, weight='bold')
    [sp.set_visible(False) for sp in ax.spines.values()]; ax.tick_params(length=0)
    savefig(path)

def fig_a1_lineage(path):
    fig, ax = plt.subplots(figsize=(7,4.2))
    ax.axis('off'); ax.set_xlim(0,10); ax.set_ylim(0,6)
    circles=[(3,3.2,1.55,navy,'Lineage 4\n(Euro-American)\npredominant'),(7,3.2,1.1,teal,'Cameroon\nsublineage\nstrongly represented'),(5,1.4,0.95,'#B0BEC5','Other MTB\nlineages also\npresent')]
    for x,y,r,c,label in circles:
        ax.add_patch(Circle((x,y),r,facecolor=c,alpha=0.9,edgecolor='white',linewidth=3))
        ax.text(x,y,label,ha='center',va='center',fontsize=11,color='white' if c!= '#B0BEC5' else navy,weight='bold')
    ax.plot([4.4,5.0],[2.4,1.9],color=gray,linewidth=2); ax.plot([6.0,5.0],[2.4,1.9],color=gray,linewidth=2)
    ax.set_title('Phylogenomic message from Abstract 1', fontsize=14, color=navy, weight='bold')
    savefig(path)

def fig_a1_diagnostic(path):
    fig, ax = plt.subplots(figsize=(7,4.2))
    ax.axis('off'); ax.set_xlim(0,10); ax.set_ylim(0,6)
    ax.add_patch(Rectangle((0.5,1.1),4.0,3.8,facecolor='#FDEDEC',edgecolor=red,linewidth=2))
    ax.add_patch(Rectangle((5.5,1.1),4.0,3.8,facecolor='#E9F7EF',edgecolor=green,linewidth=2))
    ax.text(2.5,4.4,'Canonical-only assays',ha='center',fontsize=13,weight='bold',color=red)
    ax.text(2.5,3.0,'Strong at detecting\nkatG S315T / common\ninhA promoter markers\n\nBut may miss\nnon-canonical resistance',ha='center',va='center',fontsize=11,color=navy)
    ax.text(7.5,4.4,'Expanded strategy',ha='center',fontsize=13,weight='bold',color=green)
    ax.text(7.5,3.0,'Combine WGS with MIC\n\nCaptures rare katG, inhA\nand ahpC-associated patterns\n\nImproves surveillance and\ninterpretation',ha='center',va='center',fontsize=11,color=navy)
    ax.add_patch(FancyArrowPatch((4.7,3.0),(5.3,3.0),arrowstyle='-|>',mutation_scale=18,color=teal,linewidth=2))
    ax.set_title('Why the findings matter', fontsize=14, color=navy, weight='bold')
    savefig(path)

def fig_a2_funnel(path):
    fig, ax = plt.subplots(figsize=(7,4.2))
    ax.axis('off'); ax.set_xlim(0,10); ax.set_ylim(0,8)
    widths=[8.5,6.2,4.6,3.6]; labels=['410 previously identified\nINH-resistant isolates','56 lacking canonical\nkatG/fabG1-inhA mutations','39 retained after\nphenotypic reconfirmation','30 sequenced by WGS\n(PacBio Onso)']; y=6.4; colors=['#DCEAF7','#C7DEF2','#A8CEE9','#7FB8DE']
    for w,label,c in zip(widths,labels,colors):
        x=(10-w)/2; ax.add_patch(Rectangle((x,y),w,1.1,facecolor=c,edgecolor=navy,linewidth=2)); ax.text(5,y+0.55,label,ha='center',va='center',fontsize=11,color=navy,weight='bold'); y-=1.5
    ax.text(5,0.7,'MIC results were valid for 32 isolates',ha='center',va='center',fontsize=11,color=gray)
    ax.set_title('Selection and analysis pipeline', fontsize=15, color=navy, weight='bold')
    savefig(path)

def fig_a2_mic(path):
    fig, ax = plt.subplots(figsize=(7,4.2))
    variants=['inhA\nc.-154G>A\n(n=3)', 'ahpC\nc.-74G>A\n(n=1)', 'katG\np.His400dup\n(n=3)', 'katG\np.Ser315Asn\n(n=1)', 'katG\np.Gly630fs\n(n=1)']
    mic=[0.25,0.25,0.25,8,16]; colors=[gold,gold,teal,red,red]
    ax.scatter(range(len(variants)), mic, s=[180,180,180,180,220], c=colors, edgecolor='black', linewidth=0.6)
    ax.axhline(0.25, linestyle='--', color=gray, linewidth=1.5); ax.text(4.5,0.33,'Resistance threshold = 0.25 µg/mL',fontsize=9,color=gray,ha='right')
    ax.set_yscale('log'); ax.set_ylim(0.15,20); ax.set_xticks(range(len(variants))); ax.set_xticklabels(variants, fontsize=10)
    ax.set_ylabel('MIC (µg/mL, log scale)', fontsize=11); ax.set_title('Variant-linked MIC patterns reported in Abstract 2', fontsize=14, color=navy, weight='bold'); ax.grid(axis='y', alpha=0.25)
    savefig(path)

def fig_a2_breakpoint(path):
    fig, ax = plt.subplots(figsize=(7,4.2))
    sizes=[16,16]; labels=['Below breakpoint\n16/32','At or above breakpoint\n16/32']; colors=['#BDC3C7', red]
    wedges, _ = ax.pie(sizes, startangle=90, colors=colors, wedgeprops=dict(edgecolor='white', linewidth=2))
    ax.text(0,0,'50%\nbelow\nbreakpoint',ha='center',va='center',fontsize=14,weight='bold',color=navy); ax.legend(wedges, labels, loc='center left', bbox_to_anchor=(1.0,0.5), frameon=False)
    ax.set_title('Discordance with prior MGIT classification', fontsize=14, color=navy, weight='bold')
    savefig(path)

def fig_a2_phylo(path):
    fig, ax = plt.subplots(figsize=(7,4.2))
    ax.axis('off'); ax.set_xlim(0,10); ax.set_ylim(0,6)
    ax.text(5,5.4,'Phylogenetic context from Abstract 2',ha='center',fontsize=14,weight='bold',color=navy)
    ax.plot([5,5],[4.6,3.2],color=navy,linewidth=2); ax.plot([5,2.7],[4.0,2.2],color=navy,linewidth=2); ax.plot([5,7.3],[4.0,2.2],color=navy,linewidth=2)
    ax.add_patch(Rectangle((1.2,1.1),3.0,1.1,facecolor='#DCEAF7',edgecolor=navy,linewidth=2)); ax.add_patch(Rectangle((5.8,1.1),3.0,1.1,facecolor='#EAF5FF',edgecolor=teal,linewidth=2))
    ax.text(2.7,1.65,'Lineage 4 predominant\nCameroon sublineage enriched',ha='center',va='center',fontsize=11,color=navy,weight='bold')
    ax.text(7.3,1.65,'inhA c.-154G>A\nconfined to LAM\nsublineage',ha='center',va='center',fontsize=11,color=navy,weight='bold')
    ax.text(5,3.05,'Maximum-likelihood\nphylogenetic analysis',ha='center',va='center',fontsize=11,color=gray)
    savefig(path)

assets = {'a1_fig1': imgdir/'a1_fig1_workflow.png','a1_fig2': imgdir/'a1_fig2_variants.png','a1_fig3': imgdir/'a1_fig3_lineage.png','a1_fig4': imgdir/'a1_fig4_diagnostic.png','a2_fig1': imgdir/'a2_fig1_funnel.png','a2_fig2': imgdir/'a2_fig2_mic.png','a2_fig3': imgdir/'a2_fig3_breakpoint.png','a2_fig4': imgdir/'a2_fig4_phylo.png'}
fig_a1_workflow(assets['a1_fig1']); fig_a1_variant_classes(assets['a1_fig2']); fig_a1_lineage(assets['a1_fig3']); fig_a1_diagnostic(assets['a1_fig4']); fig_a2_funnel(assets['a2_fig1']); fig_a2_mic(assets['a2_fig2']); fig_a2_breakpoint(assets['a2_fig3']); fig_a2_phylo(assets['a2_fig4'])

def clear_and_set_text(shape, paragraphs, color=RGBColor(34,52,71), center=False, line_spacing=1.0):
    tf = shape.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE; tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(4); tf.margin_bottom = Pt(4); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    for i, item in enumerate(paragraphs):
        text, level, is_bold, fs = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.text = text; p.level = level; p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT; p.space_after = Pt(2); p.space_before = Pt(0); p.line_spacing = line_spacing
        for r in p.runs:
            r.font.size = Pt(fs); r.font.color.rgb = color; r.font.bold = is_bold; r.font.name = 'Arial'

def style_title(shape, title_text):
    tf = shape.text_frame; tf.clear(); tf.word_wrap=True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run(); r.text = title_text; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = RGBColor(20,62,110); r.font.name='Arial'

def set_heading(shape, text):
    clear_and_set_text(shape, [(text,0,True,22)], color=RGBColor(20,62,110), center=False)

def set_caption(shape, text):
    clear_and_set_text(shape, [(text,0,False,11)], color=RGBColor(80,80,80), center=True)

def make_poster(template, outfile, spec):
    prs = Presentation(str(template)); slide = prs.slides[0]; placeholders = {sh.placeholder_format.idx: sh for sh in slide.placeholders if sh.is_placeholder}
    style_title(placeholders[0], spec['title'])
    clear_and_set_text(placeholders[14], spec['subtitle'], color=RGBColor(70,70,70), center=True)
    clear_and_set_text(placeholders[13], spec['topbox'])
    set_heading(placeholders[15], spec['section1_title']); clear_and_set_text(placeholders[16], spec['section1_body'])
    set_heading(placeholders[17], spec['section2_title']); set_heading(placeholders[19], spec['section3_title'])
    clear_and_set_text(placeholders[20], spec['section2_body']); clear_and_set_text(placeholders[21], spec['section3_body'])
    clear_and_set_text(placeholders[24], spec['lower_left']); clear_and_set_text(placeholders[27], spec['lower_right'])
    set_heading(placeholders[22], spec['conclusion_title']); clear_and_set_text(placeholders[23], spec['conclusion_body'])
    clear_and_set_text(placeholders[34], [(spec['footer_date'],0,False,12)], color=RGBColor(90,90,90), center=True)
    for sh in slide.shapes:
        if sh.name == 'ZoneTexte 23': clear_and_set_text(sh, [(spec['logo_text'],0,True,12)], color=RGBColor(20,62,110), center=True)
        elif sh.name == 'ZoneTexte 25': clear_and_set_text(sh, [(spec['banner_text'],0,True,16)], color=RGBColor(20,62,110), center=True)
    placeholders[25].insert_picture(str(spec['images'][0])); placeholders[26].insert_picture(str(spec['images'][1])); placeholders[28].insert_picture(str(spec['images'][2])); placeholders[29].insert_picture(str(spec['images'][3]))
    set_caption(placeholders[30], spec['captions'][0]); set_caption(placeholders[31], spec['captions'][1]); set_caption(placeholders[32], spec['captions'][2]); set_caption(placeholders[33], spec['captions'][3])
    prs.save(str(outfile))

spec1 = {'title':'Unravelling the Molecular Basis of Isoniazid Resistance in Mycobacterium tuberculosis Clinical Isolates in Cameroon with Unknown Resistance Mechanisms','subtitle':[('Vanessa Ninkeh Nono¹, Edouard Akono Nantia¹, Ndivhuho Agnes Makhado², Awelani Mutshembele³, Eric Mensah⁴, Valerie Flore Donkeng⁵',0,False,11),('¹University of Bamenda, Cameroon | ²Africa Health Research Institute, South Africa | ³South African Medical Research Council | ⁴Sefako Makgatho Health Sciences University | ⁵Centre Pasteur of Cameroon',0,False,9),('Presenter: Valerie Flore Donkeng | Email: donkeng@pasteur-yaounde.org',0,False,10)],'topbox':[('Background & study aim',0,True,18),('• Isoniazid (INH) is a core first-line TB drug, but resistance is rising and threatens TB control.',0,False,15),('• Routine rapid molecular assays focus on canonical katG codon 315 and inhA promoter mutations.',0,False,15),('• This leaves a diagnostic blind spot for phenotypically resistant isolates carrying non-canonical variants or lineage-influenced resistance patterns.',0,False,15),('• Abstract 1 examined the molecular basis of INH resistance and the genetic diversity of Cameroonian MTB isolates lacking the usual canonical markers.',0,False,15)],'section1_title':'Methodology','section1_body':[('• Focused on the 13.7% subset of previously identified INH-resistant isolates that lacked canonical katG and inhA promoter mutations.',0,False,15),('• Quantified resistance magnitude using a resazurin microtiter MIC assay.',0,False,15),('• Performed whole-genome sequencing on the PacBio® Onso platform.',0,False,15),('• Processed FASTQ data with TB-Profiler v6.6.2 for in silico resistance prediction and lineage assignment.',0,False,15),('• Compared WGS-based resistance profiles with corresponding MIC phenotypes.',0,False,15)],'section2_title':'Key molecular findings','section2_body':[('• katG-associated variants tracked with high-level INH resistance (≥8 µg/mL).',0,False,15),('• Two highlighted katG variants were c.1885_1886dupGG (p.Gly630fs) and c.944G>A (p.Ser315Asn).',0,False,15),('• The inhA promoter mutation c.-154G>A was the most frequent detected variant and was linked to low-level resistance (0.25 µg/mL).',0,False,15),('• The ahpC promoter mutation c.-74G>A was also detected and associated with low-level resistance.',0,False,15)],'section3_title':'Phylogenomics and interpretation','section3_body':[('• WGS identified four resistance-associated variants overall, alongside several missense variants of uncertain significance.',0,False,15),('• Isolates spanned diverse MTB lineages.',0,False,15),('• Lineage 4 (Euro-American) predominated, with strong representation of the Cameroon sublineage.',0,False,15),('• The abstract supports a model in which both variant type and strain background shape apparent resistance patterns.',0,False,15)],'lower_left':[('Poster highlights',0,True,17),('• Abstract 1 is framed as a diagnostic-surveillance study rather than a purely mutation cataloguing exercise.',0,False,14),('• The strongest message is that non-canonical katG, inhA, and ahpC changes can matter clinically.',0,False,14),('• Because canonical-marker-only assays are narrow, they may miss resistant isolates relevant to treatment guidance and public-health surveillance.',0,False,14),('• Keywords reported in the abstract: isoniazid resistance, Mycobacterium tuberculosis, whole-genome sequencing, resistance mechanisms, Cameroon.',0,False,13)],'lower_right':[('Why this poster matters',0,True,17),('• The work expands the local genomic picture of INH resistance in Cameroon.',0,False,14),('• It links genotype with MIC magnitude, which is more informative than a simple resistant/susceptible label.',0,False,14),('• It also argues for combining broader genomic approaches with phenotypic MIC testing to improve rapid diagnosis and TB control.',0,False,14),('• The abstract alone does not provide sample-by-sample counts beyond the 13.7% subset, so the poster intentionally stays close to what was explicitly stated.',0,False,13)],'conclusion_title':'Take-home conclusion','conclusion_body':[('• Abstract 1 shows that INH resistance in Cameroon is not limited to the classic katG S315T / common inhA pattern.',0,False,15),('• Alternative katG variants, inhA c.-154G>A, and ahpC c.-74G>A help explain resistant phenotypes in isolates that would be missed by narrower assays.',0,False,15),('• Integrating broader genomic testing with quantitative MIC measurement should strengthen diagnosis, surveillance, and treatment-oriented interpretation.',0,False,15)],'footer_date':'Prepared from Abstract 1','logo_text':'Add institutional logos here','banner_text':'Poster design based on supplied Abstract 1','images':[assets['a1_fig1'],assets['a1_fig2'],assets['a1_fig3'],assets['a1_fig4']],'captions':['Fig. 1. Workflow used in Abstract 1.','Fig. 2. Variant classes linked to low- vs high-level resistance.','Fig. 3. Lineage 4 / Cameroon sublineage message from the abstract.','Fig. 4. Why canonical-only diagnostics can miss resistant isolates.']}

spec2 = {'title':'Unravelling the Molecular Basis of Isoniazid Resistance in Mycobacterium tuberculosis Clinical Isolates in Cameroon with Unknown Resistance Mechanisms','subtitle':[('Poster concept based on supplied Abstract 2 text',0,True,11),('The Abstract 2 file did not repeat the author block; this poster therefore emphasizes the results and design logic stated in that file.',0,False,10)],'topbox':[('Background & objective',0,True,18),('• INH-resistant TB remains clinically important and is incompletely captured by rapid diagnostics centered on canonical katG codon 315 and fabG1-inhA promoter mutations.',0,False,15),('• Some phenotypically resistant isolates therefore remain genomically unexplained and can be misclassified.',0,False,15),('• Abstract 2 set out to define the genomic basis, resistance magnitude, and phylogenetic context of such isolates from Cameroon.',0,False,15)],'section1_title':'Quantitative study design','section1_body':[('• Starting set: 410 previously identified INH-resistant isolates.',0,False,15),('• 56 lacked canonical mutations; 39 remained after phenotypic reconfirmation.',0,False,15),('• MICs were measured by resazurin microtiter assay with a 0.25 µg/mL breakpoint.',0,False,15),('• WGS was completed for 30 isolates on the PacBio® Onso platform.',0,False,15),('• Variants were interpreted with TB-Profiler v6.6.2 and the WHO mutation catalogue; phylogeny used a maximum-likelihood approach.',0,False,15)],'section2_title':'Core results','section2_body':[('• Phenotypic heterogeneity was substantial, with discordance versus prior MGIT-based classification.',0,False,15),('• 16/32 isolates with valid MIC results fell below the resistance breakpoint.',0,False,15),('• High-confidence or literature-curated resistance-associated variants were found in 7/30 sequenced isolates.',0,False,15),('• inhA c.-154G>A occurred in 3 isolates and sat at MIC 0.25 µg/mL.',0,False,15),('• Rare katG variants p.Ser315Asn and p.Gly630fs were linked to MIC 8 µg/mL and ≥16 µg/mL, respectively.',0,False,15)],'section3_title':'Interpretation of non-canonical resistance','section3_body':[('• ahpC c.-74G>A was detected once at MIC 0.25 µg/mL.',0,False,15),('• One high-level resistant isolate lacked any known resistance-associated variant, reinforcing the possibility of uncharacterized mechanisms.',0,False,15),('• A novel katG p.His400dup variant appeared in 3 isolates with MICs at the resistance threshold.',0,False,15),('• Phylogeny showed predominance of Lineage 4, especially the Cameroon sublineage; inhA c.-154G>A was confined to the LAM sublineage.',0,False,15)],'lower_left':[('What changed versus a binary phenotype?',0,True,17),('• Abstract 2 makes the genotype-phenotype discordance explicit: many isolates originally labelled resistant were below the MIC breakpoint on repeat quantitative testing.',0,False,14),('• This is a key poster message because it argues against over-reliance on single-threshold phenotyping or narrow molecular panels alone.',0,False,14),('• The most data-rich result is the coupling of named variants to measured MIC values.',0,False,14)],'lower_right':[('Implications for diagnostics and surveillance',0,True,17),('• Borderline and atypical INH resistance cannot be fully resolved by canonical-marker testing alone.',0,False,14),('• Quantitative MIC plus WGS improves classification of low-level, borderline, and rare resistance patterns.',0,False,14),('• The abstract calls for functional validation of novel variants and expansion of curated mutation catalogues.',0,False,14),('• This framing is particularly relevant for TB control programs in Cameroon and similar settings.',0,False,14)],'conclusion_title':'Take-home conclusion','conclusion_body':[('• Abstract 2 presents a more quantitative view of the same diagnostic problem: INH resistance outside canonical markers is heterogeneous and only partly explained by current catalogues.',0,False,15),('• The major operational lesson is that WGS interpretation improves when it is paired with measured MICs rather than a simple resistant/susceptible label.',0,False,15),('• Further validation of novel and lineage-linked variants is needed to improve genomic resistance prediction.',0,False,15)],'footer_date':'Prepared from Abstract 2','logo_text':'Add institutional / funder logos here','banner_text':'Poster design based on supplied Abstract 2','images':[assets['a2_fig1'],assets['a2_fig2'],assets['a2_fig3'],assets['a2_fig4']],'captions':['Fig. 1. Isolate funnel and analysis pipeline.','Fig. 2. Named variants and linked MIC values from the abstract.','Fig. 3. Half of valid MIC-tested isolates fell below the breakpoint.','Fig. 4. Lineage context and sublineage restriction noted in Abstract 2.']}

make_poster(wd/'Poster template.pptx', wd/'Poster_Abstract_1_Designed.pptx', spec1)
make_poster(wd/'Poster template.pptx', wd/'Poster_Abstract_2_Designed.pptx', spec2)
print('done')
