
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Example script: write outputs beside this script by default.
# For real work, copy/adapt and use a user-supplied output directory.
WD = Path(__file__).resolve().parent
IMG = WD / "pu239_decay_chain.png"
OUT = WD / "Pu239_Decay_Chain_Brief.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NAVY = RGBColor(24, 44, 82)
BLUE = RGBColor(41, 98, 255)
TEAL = RGBColor(0, 137, 123)
GOLD = RGBColor(255, 179, 0)
LIGHT = RGBColor(245, 247, 250)
MID = RGBColor(90, 105, 125)
DARK = RGBColor(35, 45, 55)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(232, 239, 250)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def set_bg(slide, color=LIGHT):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def accent(slide):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.14))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    bar2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.6), 0, Inches(3.733), Inches(0.14))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = TEAL; bar2.line.fill.background()


def textbox(slide, text, x, y, w, h, size=20, color=DARK, bold=False, align=PP_ALIGN.LEFT, font='Aptos'):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=18, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'• {item}'
        p.font.name = 'Aptos'
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
    return box


def card(slide, x, y, w, h, fill=WHITE, line=PALE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    return shp


# Slide 1
slide = prs.slides.add_slide(blank)
set_bg(slide, LIGHT); accent(slide)
textbox(slide, 'Pu-239 Decay Chain', 0.7, 0.45, 6.0, 0.55, size=28, color=NAVY, bold=True, font='Aptos Display')
textbox(slide, 'Brief summary with actinium-series endpoint at stable Pb-207', 0.72, 0.98, 7.2, 0.3, size=15, color=MID)
card(slide, 0.75, 1.55, 5.6, 4.85, WHITE, PALE)
bullets(slide, [
    'Pu-239 decays by alpha emission to U-235.',
    'Pu-239 half-life is about 24,110 years.',
    'U-235 then follows the actinium decay series.',
    'The overall chain ends at stable Pb-207.',
    'Most steps are alpha decays, with several beta-minus transitions.',
], 1.0, 2.0, 5.1, 3.5, size=20)
card(slide, 6.7, 1.55, 5.8, 4.85, RGBColor(233,244,255), RGBColor(210,225,245))
textbox(slide, 'Main progression', 7.0, 1.85, 5.2, 0.3, size=20, color=BLUE, bold=True)
bullets(slide, [
    'Pu-239 → U-235 → Th-231 → Pa-231 → Ac-227',
    'Then to Th-227 / Ra-223 and onward through short-lived daughters',
    'Final stable product: Pb-207',
], 7.0, 2.3, 5.0, 1.9, size=17)
textbox(slide, 'Minor branching', 7.0, 4.55, 5.0, 0.25, size=18, color=TEAL, bold=True)
bullets(slide, [
    'Ac-227 has a small alpha branch to Fr-223.',
    'Bi-211 also has a minor beta-minus branch to Po-211.',
], 7.0, 4.95, 4.9, 1.2, size=16)
textbox(slide, '1', 12.35, 7.03, 0.35, 0.2, size=10, color=MID, align=PP_ALIGN.RIGHT)

# Slide 2
slide = prs.slides.add_slide(blank)
set_bg(slide, LIGHT); accent(slide)
textbox(slide, 'Decay chain diagram', 0.7, 0.45, 6.0, 0.5, size=28, color=NAVY, bold=True, font='Aptos Display')
textbox(slide, 'Mermaid-style flowchart of the main path and minor branches', 0.72, 0.98, 7.0, 0.3, size=15, color=MID)
card(slide, 0.7, 1.45, 11.95, 5.55, WHITE, PALE)
if IMG.exists():
    slide.shapes.add_picture(str(IMG), Inches(0.95), Inches(1.75), width=Inches(11.45), height=Inches(4.8))
textbox(slide, 'Alpha decay starts the chain at Pu-239; the daughter U-235 feeds into the actinium series until stable Pb-207 is reached.', 0.95, 6.32, 10.8, 0.34, size=15, color=DARK, align=PP_ALIGN.CENTER)
textbox(slide, '2', 12.35, 7.03, 0.35, 0.2, size=10, color=MID, align=PP_ALIGN.RIGHT)

# Slide 3
slide = prs.slides.add_slide(blank)
set_bg(slide, LIGHT); accent(slide)
textbox(slide, 'Key takeaways', 0.7, 0.45, 4.5, 0.5, size=28, color=NAVY, bold=True, font='Aptos Display')
textbox(slide, 'Why this chain matters in a concise technical note', 0.72, 0.98, 5.8, 0.3, size=15, color=MID)
for (x, title, body, color) in [
    (0.85, 'Decay mode', 'Pu-239 begins with alpha decay. The chain thereafter alternates between alpha and beta-minus steps.', BLUE),
    (4.45, 'Series identity', 'Once U-235 is produced, the sequence follows the actinium series, a well-known natural decay family.', TEAL),
    (8.05, 'Endpoint', 'The chain terminates at Pb-207, which is stable. Minor branches still converge to the same endpoint.', GOLD),
]:
    card(slide, x, 1.7, 3.0, 3.7, WHITE, PALE)
    textbox(slide, title, x+0.18, 2.0, 2.6, 0.28, size=18, color=color, bold=True)
    textbox(slide, body, x+0.18, 2.45, 2.55, 1.95, size=17, color=DARK)
card(slide, 1.05, 5.75, 11.2, 0.8, RGBColor(233,244,255), RGBColor(210,225,245))
textbox(slide, 'Representative half-lives: Pu-239 ≈ 24,110 y; U-235 ≈ 7.04×10^8 y. Most downstream daughters are much shorter lived.', 1.35, 6.0, 10.6, 0.25, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
textbox(slide, '3', 12.35, 7.03, 0.35, 0.2, size=10, color=MID, align=PP_ALIGN.RIGHT)

prs.save(OUT)
print(f'Saved presentation to: {OUT}')
